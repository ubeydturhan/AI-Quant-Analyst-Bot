# Auteur : Ubeyd TURHAN - 2026
from pptx import Presentation
from pptx.util import Inches, Pt
import datetime
import quant_automatisation as qa  # On importe notre cerveau Quant !
import time
from pptx.enum.text import MSO_AUTO_SIZE

print(" Création du Rapport PowerPoint Hebdomadaire...")

# 1. On prépare le document
prs = Presentation("template_quant.pptx") # Ça, ça ouvre ton fichier avec tout ton design !

# 2. La Slide de Titre
slide_titre = prs.slides.add_slide(prs.slide_layouts[0]) # Layout 0 = Titre principal
titre = slide_titre.shapes.title
sous_titre = slide_titre.placeholders[1]

titre.text = "Rapport Quantitatif IA"
sous_titre.text = f"Généré automatiquement le {datetime.date.today()}\nAnalyse par Agent Kimi"

# 3. On analyse notre portefeuille en silence et on crée les slides
mon_portefeuille = ["MSFT", "AAPL", "NVDA"]

for ticker in mon_portefeuille:
    print(f" Traitement de {ticker} en cours pour le PowerPoint...")
    
    # On fait appel à nos fonctions de l'autre fichier !
    prix = qa.analyser_prix(ticker)
    verdict_ia = qa.analyser_action(ticker, prix)
    
    # On ajoute une nouvelle Slide (Layout 1 = Titre + Contenu)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    # On remplit le titre de la slide
    slide.shapes.title.text = f"Analyse Stratégique : {ticker}"
    
    # On remplit le corps de la slide avec l'analyse de Kimi
    corps_texte = slide.placeholders[1].text_frame
    corps_texte.word_wrap = True  # Force le texte à aller à la ligne
    corps_texte.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE # Demande à PPT de réduire la taille
    
    corps_texte.text = verdict_ia
    
    # On réduit un peu la police pour que tout rentre sur la page
    for paragraph in corps_texte.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(11) 
            
    print(f"✅ Slide pour {ticker} ajoutée !")
    time.sleep(12) # Pause anti-surchauffe pour Kimi

# 4. On sauvegarde le fichier physique !
nom_fichier = "Rapport_Quant.pptx"
prs.save(nom_fichier)

print(f"\n🎉 SUCCÈS TOTAL ! Ouvre le fichier '{nom_fichier}' dans ton dossier !")